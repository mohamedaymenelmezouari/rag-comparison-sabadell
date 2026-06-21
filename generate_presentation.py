"""
generate_presentation.py — Final-submission PowerPoint for ESADE Capstone 2026.
Run: python3 generate_presentation.py
Requires: python-pptx, Pillow
"""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Colours ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x07, 0x0D, 0x1A)   # near-black navy
CARD    = RGBColor(0x0C, 0x15, 0x26)   # card bg
CARD2   = RGBColor(0x0F, 0x1C, 0x35)   # slightly lighter card
NAVY    = RGBColor(0x00, 0x20, 0x54)   # Banc Sabadell navy
BLUE    = RGBColor(0x00, 0x6D, 0xFF)   # Banc Sabadell blue
BLUE_L  = RGBColor(0x3B, 0x9E, 0xFF)   # lighter blue
TRAD    = RGBColor(0x37, 0x99, 0xF6)   # Traditional RAG
GRAPH   = RGBColor(0x10, 0xB9, 0x81)   # Graph RAG green
GRAPH_D = RGBColor(0x05, 0x96, 0x69)   # darker green
GOLD    = RGBColor(0xF5, 0x9E, 0x0B)   # amber
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGREY   = RGBColor(0x94, 0xA3, 0xB8)
DGREY   = RGBColor(0x33, 0x41, 0x55)
BORDER  = RGBColor(0x13, 0x1F, 0x3A)
RED_ACC = RGBColor(0xEF, 0x44, 0x44)

LOGO_PATH = Path("/tmp/sabadell_logo_dark.png")
OUT_PATH  = Path("/Users/yeezyuus/Downloads/Capstone_Project_Banc_Sabadell (2)/Capstone Project Banc Sabadell/RAG_vs_GraphRAG_Final_Presentation.pptx")

W = Inches(13.33)
H = Inches(7.5)


# ── Core helpers ──────────────────────────────────────────────────────────────
def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    bg     = slide.background
    fill   = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def rect(slide, l, t, w, h, fill=None, line=None, lw=0.5, radius=False):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.05
    return shape


def rrect(slide, l, t, w, h, fill=None, line=None, lw=0.5):
    """Rounded rectangle."""
    shape = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, l, t, w, h, text="", size=12, bold=False,
            color=WHITE, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def multitext(slide, l, t, w, h, lines: list[tuple], wrap=True):
    """lines = [(text, size, bold, color, align, space_before), ...]"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    for i, (text, size, bold, color, align, sp) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if sp:
            p.space_before = Pt(sp)
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
    return txb


def logo(slide, l, t, w=2.2):
    if LOGO_PATH.exists():
        return slide.shapes.add_picture(str(LOGO_PATH), Inches(l), Inches(t), Inches(w))


def accent_bar(slide, color=BLUE, height=0.055):
    """Full-width coloured bar at very top."""
    r = rect(slide, 0, 0, 13.33, height, fill=color)
    return r


def slide_header(slide, title_text: str, show_logo=True):
    accent_bar(slide)
    textbox(slide, 0.45, 0.12, 9.5, 0.55,
            title_text, size=22, bold=True, color=WHITE)
    if show_logo:
        logo(slide, 10.9, 0.06, 2.1)


def add_fade_transition(slide):
    """Add fade transition to a slide via raw XML."""
    trans_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" dur="400"><p:fade/></p:transition>'
    elem = parse_xml(trans_xml)
    slide._element.append(elem)


def circle(slide, l, t, d, fill):
    """Add a circle (oval) shape."""
    shape = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def arrow_right(slide, l, t, w=0.25, h=0.12):
    """Small right arrow connector."""
    shape = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DGREY
    shape.line.fill.background()
    return shape


def divider(slide, l, t, w, color=BORDER, vertical=False):
    if vertical:
        r = rect(slide, l, t, 0.005, w, fill=color)
    else:
        r = rect(slide, l, t, w, 0.005, fill=color)
    return r


# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
def s_title(prs):
    slide = blank_slide(prs)

    # Subtle glow orbs
    glow = slide.shapes.add_shape(9, Inches(4.2), Inches(1.5), Inches(5), Inches(5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x00, 0x25, 0x60)
    glow.line.fill.background()
    # make it semi-transparent via xml
    sp_pr = glow._element.find(qn('p:spPr'))
    solid_fill = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solid_fill is not None:
        srgb = solid_fill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '18000')  # 18% opacity

    glow2 = slide.shapes.add_shape(9, Inches(8.5), Inches(2.5), Inches(3.5), Inches(3.5))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = RGBColor(0x00, 0x40, 0x90)
    glow2.line.fill.background()
    sp_pr2 = glow2._element.find(qn('p:spPr'))
    solid_fill2 = sp_pr2.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solid_fill2 is not None:
        srgb2 = solid_fill2.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb2 is not None:
            alpha2 = etree.SubElement(srgb2, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha2.set('val', '12000')

    # Top blue accent bar
    rect(slide, 0, 0, 13.33, 0.065, fill=BLUE)

    # Logo
    logo(slide, 5.3, 0.55, 2.72)

    # Thin divider below logo
    divider(slide, 3.5, 1.25, 6.33, color=RGBColor(0x1E, 0x3A, 0x6E))

    # Main title
    multitext(slide, 1.0, 1.45, 11.33, 2.2, [
        ("Traditional RAG  vs.  Graph RAG", 46, True, WHITE, PP_ALIGN.CENTER, 0),
        ("Banking Regulatory Intelligence for Banc Sabadell", 20, False, BLUE_L, PP_ALIGN.CENTER, 14),
    ])

    # Bottom stripe
    rect(slide, 0, 6.35, 13.33, 1.15, fill=RGBColor(0x03, 0x08, 0x12))
    divider(slide, 0, 6.35, 13.33, color=BLUE)

    # Team and context
    multitext(slide, 0.5, 6.5, 12.33, 0.7, [
        ("Mohamed Aymen Elmezouari  ·  Abhay Mathahalli  ·  Warren Liu",
         11, False, LGREY, PP_ALIGN.CENTER, 0),
        ("ESADE Capstone 2026  |  Final Presentation  |  June 2026",
         10, False, DGREY, PP_ALIGN.CENTER, 4),
    ])

    add_fade_transition(slide)
    return slide


# ── SLIDE 2: AGENDA ───────────────────────────────────────────────────────────
def s_agenda(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Agenda")

    items = [
        ("01", "The Business Challenge",              BLUE),
        ("02", "What is RAG?  Traditional vs. Graph", BLUE),
        ("03", "System Architecture",                 BLUE),
        ("04", "Benchmark Results (75 queries)",      GOLD),
        ("05", "The Key Insight: Accuracy by Tier",   GOLD),
        ("06", "Recommendation: Hybrid Architecture", GRAPH),
        ("07", "GDPR Compliance  +  Roadmap",         GRAPH),
    ]

    col_w = 5.8
    for i, (num, label, color) in enumerate(items):
        row = i % 4
        col = i // 4
        lx  = 0.45 + col * 6.4
        ty  = 0.95 + row * 1.52

        rrect(slide, lx, ty, col_w, 1.22, fill=CARD, line=BORDER, lw=0.5)
        rect(slide, lx, ty, 0.055, 1.22, fill=color)

        textbox(slide, lx + 0.18, ty + 0.08, 0.65, 0.45,
                num, size=26, bold=True, color=color, align=PP_ALIGN.LEFT)
        textbox(slide, lx + 0.78, ty + 0.22, col_w - 0.95, 0.65,
                label, size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 7th item spans differently
    rrect(slide, 0.45, 0.95 + 3 * 1.52, col_w, 1.22, fill=CARD, line=BORDER, lw=0.5)
    rect(slide, 0.45, 0.95 + 3 * 1.52, 0.055, 1.22, fill=GRAPH)
    textbox(slide, 0.63, 0.95 + 3 * 1.52 + 0.08, 0.65, 0.45,
            "07", size=26, bold=True, color=GRAPH, align=PP_ALIGN.LEFT)
    textbox(slide, 1.23, 0.95 + 3 * 1.52 + 0.22, col_w - 0.95, 0.65,
            "GDPR Compliance  +  Roadmap", size=14, bold=True, color=WHITE)

    add_fade_transition(slide)
    return slide


# ── SLIDE 3: THE CHALLENGE ────────────────────────────────────────────────────
def s_challenge(prs):
    slide = blank_slide(prs)
    slide_header(slide, "The Business Challenge")

    textbox(slide, 0.45, 0.82, 12.43, 0.5,
            "Banks manage thousands of pages of regulatory documents. Current systems fail on three critical fronts.",
            size=13, color=LGREY, align=PP_ALIGN.LEFT)

    cards = [
        (TRAD,  "Multi-hop Queries",
         "Answering 'which products require both MiFID II suitability AND AML enhanced due diligence?' "
         "requires linking entities across frameworks. Vector search cannot do this reliably."),
        (GOLD,  "Hallucination Risk",
         "A fabricated regulatory citation in a compliance workflow constitutes a material risk. "
         "Systems that hallucinate answers are unacceptable in a banking context."),
        (GRAPH, "Audit Trails",
         "Regulators and compliance officers must be able to trace exactly how an answer was constructed. "
         "Black-box retrieval fails MiFID II explainability requirements."),
    ]

    for i, (color, title, body_text) in enumerate(cards):
        lx = 0.45 + i * 4.3
        rrect(slide, lx, 1.55, 4.0, 4.55, fill=CARD, line=BORDER)
        rect(slide, lx, 1.55, 4.0, 0.07, fill=color)

        c = circle(slide, lx + 1.5, 1.85, 1.0, fill=RGBColor(
            int(color[0]*0.2), int(color[1]*0.2), int(color[2]*0.2)))
        textbox(slide, lx + 1.5, 1.85, 1.0, 1.0,
                str(i + 1), size=34, bold=True, color=color, align=PP_ALIGN.CENTER)

        textbox(slide, lx + 0.18, 3.1, 3.65, 0.5,
                title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, lx + 0.2, 3.65, 3.62, 2.2,
                body_text, size=10.5, color=LGREY, align=PP_ALIGN.LEFT)

    add_fade_transition(slide)
    return slide


# ── SLIDE 4: RAG EXPLAINER ────────────────────────────────────────────────────
def s_rag(prs):
    slide = blank_slide(prs)
    slide_header(slide, "What is RAG?  Why Does Architecture Matter?")

    textbox(slide, 0.45, 0.82, 12.43, 0.42,
            "Retrieval-Augmented Generation grounds language model answers in real documents, "
            "eliminating hallucination by design. The architecture of the retrieval step determines "
            "everything.",
            size=12, color=LGREY)

    # Simple RAG flow at top
    steps = ["Documents", "Retrieve", "Augment", "Generate", "Answer"]
    colors_flow = [DGREY, TRAD, TRAD, TRAD, GRAPH]
    for i, (step, c) in enumerate(zip(steps, colors_flow)):
        lx = 0.45 + i * 2.52
        rrect(slide, lx, 1.5, 2.1, 0.75, fill=RGBColor(
            max(c[0]-40, 0), max(c[1]-40, 0), max(c[2]-40, 0)), line=c, lw=1.2)
        textbox(slide, lx, 1.5, 2.1, 0.75,
                step, size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        if i < 4:
            textbox(slide, lx + 2.1, 1.73, 0.4, 0.3,
                    "▶", size=11, color=DGREY, align=PP_ALIGN.CENTER)

    divider(slide, 0.45, 2.55, 12.43, color=BORDER)

    # Two paradigm cards
    for col, (title, subtitle, color, points) in enumerate([
        ("Traditional RAG",
         "Semantic Similarity · Vector Retrieval",
         TRAD,
         ["Embeds every chunk as a vector (all-MiniLM-L6-v2)",
          "Retrieves top-k chunks by cosine similarity",
          "Passes chunks as flat context to the LLM",
          "Fast and simple to build and update",
          "Cannot reason across entity relationships"]),
        ("Graph RAG",
         "Structured Relationships · Graph Traversal",
         GRAPH,
         ["Extracts entities and relationships from documents",
          "Builds a typed knowledge graph (Neo4j)",
          "Traverses the graph with Cypher queries at query time",
          "Provides auditable multi-hop reasoning paths",
          "Significantly better on complex relational queries"]),
    ]):
        lx = 0.45 + col * 6.5
        rrect(slide, lx, 2.75, 6.15, 4.35, fill=CARD, line=BORDER)
        rect(slide, lx, 2.75, 6.15, 0.065, fill=color)
        textbox(slide, lx + 0.2, 2.88, 5.75, 0.5,
                title, size=16, bold=True, color=color)
        textbox(slide, lx + 0.2, 3.3, 5.75, 0.35,
                subtitle, size=9.5, color=DGREY)
        for j, pt in enumerate(points):
            textbox(slide, lx + 0.45, 3.72 + j * 0.62, 5.5, 0.5,
                    f"• {pt}", size=10.5, color=LGREY)

    add_fade_transition(slide)
    return slide


# ── SLIDE 5: ARCHITECTURE — TRADITIONAL RAG ───────────────────────────────────
def s_trad_arch(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Architecture: Traditional RAG Pipeline")

    steps = [
        ("PDF / TXT", "Source\nDocuments"),
        ("PyMuPDF", "PDF\nParsing"),
        ("Chunking\n500 tok", "Text\nSplitting"),
        ("MiniLM\nEmbed", "Local\nEmbedding"),
        ("ChromaDB", "Vector\nIndex"),
        ("Top-k\nCosine", "Retrieval"),
        ("Groq\nLlama 3.3", "Generation"),
    ]

    box_w, box_h = 1.52, 1.0
    start_x, y_top = 0.32, 1.7
    gap = 0.2

    for i, (main, sub) in enumerate(steps):
        lx = start_x + i * (box_w + gap)
        c  = DGREY if i == 0 else (GOLD if i == 6 else TRAD)
        fill_c = RGBColor(max(c[0]-60, 0), max(c[1]-60, 0), max(c[2]-60, 0))
        rrect(slide, lx, y_top, box_w, box_h, fill=fill_c, line=c, lw=1.5)
        textbox(slide, lx, y_top + 0.08, box_w, 0.55,
                main, size=11, bold=True, color=c, align=PP_ALIGN.CENTER)
        textbox(slide, lx, y_top + 0.6, box_w, 0.35,
                sub, size=8, color=DGREY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            textbox(slide, lx + box_w + 0.02, y_top + 0.38, gap, 0.25,
                    "▶", size=12, color=DGREY, align=PP_ALIGN.CENTER)

    # Design decisions boxes
    decisions = [
        ("ChromaDB (local)", "Full data locality — no regulatory text leaves the environment (GDPR)"),
        ("all-MiniLM-L6-v2", "Zero API cost, open-source, strong domain-general semantic performance"),
        ("500-token chunks", "Calibrated to preserve full articles/sub-sections — maximises retrieval precision"),
        ("Groq / Llama 3.3", "Sub-2-second latency at 70B parameter scale — primary performance target"),
    ]

    textbox(slide, 0.45, 3.05, 6.5, 0.38,
            "Key Design Decisions", size=11, bold=True, color=LGREY)

    for i, (label, desc) in enumerate(decisions):
        row, col = i % 2, i // 2
        lx = 0.45 + col * 6.5
        ty = 3.45 + row * 1.55
        rrect(slide, lx, ty, 6.15, 1.32, fill=CARD, line=BORDER)
        rect(slide, lx, ty, 0.055, 1.32, fill=TRAD)
        textbox(slide, lx + 0.18, ty + 0.1, 5.8, 0.38,
                label, size=12, bold=True, color=TRAD)
        textbox(slide, lx + 0.18, ty + 0.5, 5.8, 0.75,
                desc, size=10, color=LGREY)

    add_fade_transition(slide)
    return slide


# ── SLIDE 6: ARCHITECTURE — GRAPH RAG ────────────────────────────────────────
def s_graph_arch(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Architecture: Graph RAG Pipeline")

    steps = [
        ("PDF / TXT", "Source"),
        ("spaCy NER", "Entity\nExtract"),
        ("8 Entity\nTypes", "Schema"),
        ("Neo4j\nGraph", "Storage"),
        ("Cypher\nQuery", "Traversal"),
        ("Subgraph\nContext", "Assembly"),
        ("Groq\nLlama 3.3", "Generation"),
    ]

    box_w, box_h = 1.52, 1.0
    start_x, y_top = 0.32, 1.65

    for i, (main, sub) in enumerate(steps):
        lx = start_x + i * 1.74
        c  = DGREY if i == 0 else (GOLD if i == 6 else GRAPH)
        fill_c = RGBColor(max(c[0]-60, 0), max(c[1]-60, 0), max(c[2]-60, 0))
        rrect(slide, lx, y_top, box_w, box_h, fill=fill_c, line=c, lw=1.5)
        textbox(slide, lx, y_top + 0.08, box_w, 0.55,
                main, size=11, bold=True, color=c, align=PP_ALIGN.CENTER)
        textbox(slide, lx, y_top + 0.62, box_w, 0.3,
                sub, size=8, color=DGREY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            textbox(slide, lx + box_w + 0.02, y_top + 0.38, 0.2, 0.25,
                    "▶", size=12, color=DGREY, align=PP_ALIGN.CENTER)

    # Schema section
    textbox(slide, 0.45, 3.0, 12.43, 0.38,
            "Knowledge Graph Schema", size=11, bold=True, color=LGREY)

    entity_types = ["Regulation", "Article", "Product", "Obligation",
                    "Role", "Risk Category", "Threshold", "Date"]
    rel_types = ["GOVERNS", "REQUIRES", "APPLIES TO",
                 "REFERENCES", "SUPERSEDES", "EXEMPTS"]

    rrect(slide, 0.45, 3.4, 7.8, 3.7, fill=CARD, line=BORDER)
    rect(slide, 0.45, 3.4, 7.8, 0.06, fill=GRAPH)
    textbox(slide, 0.65, 3.5, 3, 0.4,
            "Entity Types (8)", size=10, bold=True, color=GRAPH)
    for i, e in enumerate(entity_types):
        row, col = i % 4, i // 4
        lx2 = 0.65 + col * 3.75
        ty2 = 3.98 + row * 0.72
        rrect(slide, lx2, ty2, 3.5, 0.52, fill=RGBColor(0x05, 0x20, 0x18), line=GRAPH_D, lw=0.8)
        textbox(slide, lx2, ty2, 3.5, 0.52,
                e, size=10, bold=False, color=GRAPH, align=PP_ALIGN.CENTER)

    rrect(slide, 8.48, 3.4, 4.4, 3.7, fill=CARD, line=BORDER)
    rect(slide, 8.48, 3.4, 4.4, 0.06, fill=GRAPH_D)
    textbox(slide, 8.68, 3.5, 4.0, 0.4,
            "Relationship Types (6)", size=10, bold=True, color=GRAPH)
    for i, r_name in enumerate(rel_types):
        ty2 = 3.98 + i * 0.53
        rrect(slide, 8.68, ty2, 3.9, 0.42, fill=RGBColor(0x03, 0x14, 0x10), line=GRAPH_D, lw=0.8)
        textbox(slide, 8.68, ty2, 3.9, 0.42,
                r_name, size=10, bold=True, color=GRAPH, align=PP_ALIGN.CENTER)

    # Key improvement callout
    rrect(slide, 0.45, 6.8, 12.43, 0.5, fill=RGBColor(0x05, 0x20, 0x18), line=GRAPH, lw=1.0)
    textbox(slide, 0.65, 6.85, 12.0, 0.4,
            "Mid-term to final: relationship density increased from 22 to 847 (38.5x improvement), "
            "enabling genuine multi-hop reasoning across the corpus.",
            size=10, color=GRAPH)

    add_fade_transition(slide)
    return slide


# ── SLIDE 7: THE NUMBERS (WOW SLIDE) ─────────────────────────────────────────
def s_wow(prs):
    slide = blank_slide(prs)

    # Large subtle glow behind stats
    glow = slide.shapes.add_shape(9, Inches(2.5), Inches(0.5), Inches(8.5), Inches(6.8))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x00, 0x30, 0x80)
    glow.line.fill.background()
    sp_pr = glow._element.find(qn('p:spPr'))
    sf = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if sf is not None:
        srgb = sf.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            a = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            a.set('val', '8000')

    rect(slide, 0, 0, 13.33, 0.065, fill=BLUE)

    textbox(slide, 0.45, 0.18, 12.43, 0.65,
            "The Results", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.45, 0.78, 12.43, 0.4,
            "75 annotated query pairs  ·  June 2026  ·  GPT-4o independent judge",
            size=12, color=DGREY, align=PP_ALIGN.CENTER)

    stats = [
        ("+21pp",  "accuracy advantage",          "on multi-hop relational queries (Tier 3)", BLUE,  BLUE_L),
        ("57%",    "reduction in hallucination",   "Graph RAG 9% vs Traditional RAG 21%",      GRAPH, GRAPH),
        ("52/75",  "queries won by Graph RAG",     "69% win rate across all complexity tiers",  GOLD,  GOLD),
        ("2.3s",   "Traditional RAG latency",      "Graph RAG adds only 0.8s overhead",         TRAD,  TRAD),
    ]

    card_w, card_h = 2.95, 4.6
    for i, (number, label1, label2, accent, txt_color) in enumerate(stats):
        lx = 0.45 + i * 3.22
        rrect(slide, lx, 1.35, card_w, card_h, fill=CARD, line=BORDER)
        rect(slide, lx, 1.35, card_w, 0.07, fill=accent)
        # Giant number
        textbox(slide, lx + 0.05, 1.65, card_w - 0.1, 1.55,
                number, size=64 if len(number) <= 4 else 48,
                bold=True, color=txt_color, align=PP_ALIGN.CENTER)
        # Label line 1
        textbox(slide, lx + 0.1, 3.38, card_w - 0.2, 0.55,
                label1, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Label line 2
        textbox(slide, lx + 0.1, 3.92, card_w - 0.2, 0.82,
                label2, size=10.5, color=LGREY, align=PP_ALIGN.CENTER)

    add_fade_transition(slide)
    return slide


# ── SLIDE 8: FULL RESULTS TABLE ───────────────────────────────────────────────
def s_results(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Full Benchmark Results  (n=75, GPT-4o Judge)")

    rows = [
        ("Accuracy",               "79%", "87%", "+8 pp"),
        ("Faithfulness",           "82%", "91%", "+9 pp"),
        ("Source Citation Quality","74%", "88%", "+14 pp"),
        ("Answer Quality",         "81%", "89%", "+8 pp"),
        ("Hallucination Risk",     "21%",  "9%", "−57%"),
        ("Average Latency",       "2.3s", "3.1s", "+0.8s"),
    ]

    headers = ["Metric", "Traditional RAG", "Graph RAG", "Delta"]
    col_widths = [3.8, 2.8, 2.8, 2.0]
    col_starts = [0.45, 4.25, 7.05, 9.85]
    row_h = 0.68

    # Header row
    for j, (hdr, lx, cw) in enumerate(zip(headers, col_starts, col_widths)):
        rrect(slide, lx, 0.95, cw, 0.5, fill=NAVY, line=BLUE, lw=0.8)
        textbox(slide, lx, 0.95, cw, 0.5,
                hdr, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (metric, trad_v, graph_v, delta) in enumerate(rows):
        ty = 1.52 + i * row_h
        bg = CARD if i % 2 == 0 else RGBColor(0x0A, 0x12, 0x22)
        is_hall = "Hallucination" in metric
        is_lat  = "Latency" in metric

        for j, (val, lx, cw) in enumerate(zip(
            [metric, trad_v, graph_v, delta], col_starts, col_widths
        )):
            rrect(slide, lx, ty, cw, row_h - 0.05, fill=bg, line=BORDER, lw=0.3)
            if j == 0:
                textbox(slide, lx + 0.15, ty, cw - 0.15, row_h,
                        val, size=11, bold=True, color=LGREY, align=PP_ALIGN.LEFT)
            elif j == 1:
                textbox(slide, lx, ty, cw, row_h,
                        val, size=15, bold=True, color=TRAD, align=PP_ALIGN.CENTER)
            elif j == 2:
                g_better = not is_lat
                textbox(slide, lx, ty, cw, row_h,
                        val, size=15, bold=True,
                        color=GRAPH if g_better else RED_ACC, align=PP_ALIGN.CENTER)
            else:
                d_color = GRAPH if (not is_lat and "+" in delta) or "−" in delta else RED_ACC
                if is_lat:
                    d_color = LGREY
                textbox(slide, lx, ty, cw, row_h,
                        delta, size=14, bold=True, color=d_color, align=PP_ALIGN.CENTER)

    # Summary bar at bottom
    rrect(slide, 0.45, 6.35, 12.43, 0.82, fill=RGBColor(0x03, 0x12, 0x2A), line=BLUE, lw=1.0)
    multitext(slide, 0.7, 6.42, 12.0, 0.7, [
        ("Overall: Graph RAG wins 52/75 queries (69%)   ·   Traditional RAG wins 18/75 (24%)   ·   5 ties (7%)",
         12, True, WHITE, PP_ALIGN.CENTER, 0),
    ])

    add_fade_transition(slide)
    return slide


# ── SLIDE 9: ACCURACY BY TIER (bar chart) ─────────────────────────────────────
def s_by_tier(prs):
    slide = blank_slide(prs)
    slide_header(slide, "The Key Insight: Accuracy Gap Grows with Query Complexity")

    textbox(slide, 0.45, 0.82, 12.43, 0.42,
            "Both architectures start at near-parity on simple factual queries. "
            "The gap widens dramatically as queries require multi-hop reasoning.",
            size=12, color=LGREY)

    # Chart area
    chart_l, chart_t = 0.45, 1.4
    chart_w, chart_h = 8.5, 5.2
    max_val = 100.0
    bar_area_h = 4.0
    bar_area_t = chart_t + 0.5

    rrect(slide, chart_l, chart_t, chart_w, chart_h, fill=CARD, line=BORDER)

    tiers = [
        ("Tier 1\nSimple Factual", 91, 89),
        ("Tier 2\nSingle-hop",     84, 90),
        ("Tier 3\nMulti-hop",      68, 89),
        ("Tier 4\nAdversarial",    71, 80),
    ]

    group_w  = 1.75
    bar_w    = 0.6
    group_gap = 0.32
    start_x  = chart_l + 0.55

    # Y axis labels
    for pct in [0, 25, 50, 75, 100]:
        y_pos = bar_area_t + bar_area_h - (pct / max_val) * bar_area_h
        divider(slide, chart_l + 0.45, y_pos, chart_w - 0.6, color=BORDER)
        textbox(slide, chart_l + 0.05, y_pos - 0.15, 0.4, 0.3,
                f"{pct}%", size=8, color=DGREY, align=PP_ALIGN.RIGHT)

    for i, (tier_label, trad_pct, graph_pct) in enumerate(tiers):
        gx = start_x + i * (group_w + group_gap)

        # Traditional RAG bar
        bh = (trad_pct / max_val) * bar_area_h
        by = bar_area_t + bar_area_h - bh
        rrect(slide, gx, by, bar_w, bh, fill=TRAD, line=None)
        textbox(slide, gx, by - 0.32, bar_w, 0.3,
                f"{trad_pct}%", size=10, bold=True, color=TRAD, align=PP_ALIGN.CENTER)

        # Graph RAG bar
        bh2 = (graph_pct / max_val) * bar_area_h
        by2 = bar_area_t + bar_area_h - bh2
        rrect(slide, gx + bar_w + 0.08, by2, bar_w, bh2, fill=GRAPH, line=None)
        textbox(slide, gx + bar_w + 0.08, by2 - 0.32, bar_w, 0.3,
                f"{graph_pct}%", size=10, bold=True, color=GRAPH, align=PP_ALIGN.CENTER)

        # Tier label
        textbox(slide, gx - 0.12, bar_area_t + bar_area_h + 0.08, group_w, 0.55,
                tier_label, size=9, color=LGREY, align=PP_ALIGN.CENTER)

        # Delta label on Tier 3 (most important)
        if i == 2:
            delta_y = min(by, by2) - 0.5
            rrect(slide, gx - 0.05, delta_y, group_w + 0.1, 0.38,
                  fill=RGBColor(0x05, 0x20, 0x18), line=GRAPH, lw=1.0)
            textbox(slide, gx - 0.05, delta_y, group_w + 0.1, 0.38,
                    "+21pp", size=13, bold=True, color=GRAPH, align=PP_ALIGN.CENTER)

    # Legend
    rrect(slide, chart_l + 0.55, chart_t + chart_h - 0.45, 1.6, 0.32,
          fill=TRAD, line=None)
    textbox(slide, chart_l + 2.2, chart_t + chart_h - 0.48, 1.5, 0.38,
            "Traditional RAG", size=9, color=TRAD)
    rrect(slide, chart_l + 3.75, chart_t + chart_h - 0.45, 1.6, 0.32,
          fill=GRAPH, line=None)
    textbox(slide, chart_l + 5.4, chart_t + chart_h - 0.48, 1.5, 0.38,
            "Graph RAG", size=9, color=GRAPH)

    # Right panel — key insights
    rrect(slide, 9.2, 1.4, 3.7, 5.2, fill=CARD, line=BORDER)
    rect(slide, 9.2, 1.4, 3.7, 0.065, fill=GOLD)
    textbox(slide, 9.4, 1.55, 3.3, 0.45,
            "Key Insights", size=13, bold=True, color=GOLD)

    insights = [
        ("Tier 1", "Near-parity — both work well for simple factual lookups", TRAD),
        ("Tier 2", "+6pp advantage as single relationships emerge", GRAPH),
        ("Tier 3", "+21pp as multi-hop reasoning becomes critical", GRAPH),
        ("Tier 4", "+9pp — Graph RAG resists adversarial hallucination better", GRAPH),
    ]
    for j, (tier, text, color) in enumerate(insights):
        ty = 2.1 + j * 1.1
        rect(slide, 9.4, ty, 0.05, 0.82, fill=color)
        textbox(slide, 9.55, ty, 3.2, 0.35,
                tier, size=10, bold=True, color=color)
        textbox(slide, 9.55, ty + 0.33, 3.2, 0.55,
                text, size=9.5, color=LGREY)

    add_fade_transition(slide)
    return slide


# ── SLIDE 10: HEAD TO HEAD ────────────────────────────────────────────────────
def s_head2head(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Head to Head: When Each Architecture Wins")

    cols = [
        ("Traditional RAG", TRAD, [
            "Tier 1 simple factual queries (91% accuracy, near-parity with Graph RAG)",
            "Latency-critical workflows: 2.3s average vs 3.1s for Graph RAG",
            "Frequently updated corpora where graph maintenance adds operational cost",
            "High-volume customer FAQ and product information lookups",
            "Infrastructure-constrained environments without a Neo4j instance",
        ]),
        ("Graph RAG", GRAPH, [
            "Tier 3 multi-hop relational queries: +21 percentage point accuracy advantage",
            "Compliance-critical workflows where hallucination risk must be minimised (9% vs 21%)",
            "MiFID II explainability requirements: full traversal path logged per query",
            "Cross-framework queries spanning Basel III, MiFID II, GDPR, and AML simultaneously",
            "Regulatory audit preparation where source citation quality (88% vs 74%) is essential",
        ]),
    ]

    for i, (title, color, points) in enumerate(cols):
        lx = 0.45 + i * 6.5
        rrect(slide, lx, 0.9, 6.1, 5.85, fill=CARD, line=BORDER)
        rect(slide, lx, 0.9, 6.1, 0.08, fill=color)

        textbox(slide, lx + 0.2, 1.05, 5.7, 0.52,
                title, size=18, bold=True, color=color)

        # Score badges
        score_lbl = "Wins 18/75 queries" if i == 0 else "Wins 52/75 queries"
        rrect(slide, lx + 0.2, 1.62, 3.2, 0.4,
              fill=RGBColor(max(color[0]-80,0), max(color[1]-80,0), max(color[2]-80,0)),
              line=color, lw=0.8)
        textbox(slide, lx + 0.2, 1.62, 3.2, 0.4,
                score_lbl, size=9.5, bold=True, color=color, align=PP_ALIGN.CENTER)

        for j, pt in enumerate(points):
            ty = 2.18 + j * 0.9
            rrect(slide, lx + 0.2, ty, 5.7, 0.75, fill=RGBColor(0x07, 0x0D, 0x1A), line=BORDER)
            rect(slide, lx + 0.2, ty, 0.045, 0.75, fill=color)
            textbox(slide, lx + 0.35, ty + 0.08, 5.4, 0.6,
                    pt, size=10, color=LGREY)

    add_fade_transition(slide)
    return slide


# ── SLIDE 11: RECOMMENDATION ──────────────────────────────────────────────────
def s_recommendation(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Our Recommendation: Hybrid Query-Routing Architecture")

    textbox(slide, 0.45, 0.82, 12.43, 0.42,
            "Neither architecture dominates across all query types. "
            "The optimal solution routes each query to the right pipeline based on complexity.",
            size=12, color=LGREY)

    # Central routing diagram
    # Query input
    rrect(slide, 0.4, 2.2, 2.1, 0.85, fill=CARD, line=DGREY, lw=1.2)
    textbox(slide, 0.4, 2.2, 2.1, 0.85,
            "Query", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow right
    textbox(slide, 2.5, 2.42, 0.5, 0.4, "▶", size=14, color=DGREY, align=PP_ALIGN.CENTER)

    # Classifier box
    rrect(slide, 3.0, 1.95, 2.9, 1.35, fill=CARD, line=GOLD, lw=1.5)
    rect(slide, 3.0, 1.95, 2.9, 0.07, fill=GOLD)
    textbox(slide, 3.0, 2.08, 2.9, 0.45,
            "Query Classifier", size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    textbox(slide, 3.1, 2.52, 2.7, 0.65,
            "BERT fine-tuned\n<100ms inference", size=9.5, color=LGREY, align=PP_ALIGN.CENTER)

    # Arrows to both pipelines
    textbox(slide, 5.9, 1.65, 0.55, 0.4, "▲ Tier 2-4", size=8.5, color=GRAPH, align=PP_ALIGN.CENTER)
    textbox(slide, 5.9, 2.88, 0.55, 0.4, "▼ Tier 1", size=8.5, color=TRAD, align=PP_ALIGN.CENTER)

    # Graph RAG path (top)
    rrect(slide, 6.55, 1.2, 3.2, 1.15, fill=RGBColor(0x03, 0x18, 0x10), line=GRAPH, lw=1.5)
    textbox(slide, 6.55, 1.26, 3.2, 0.45,
            "Graph RAG", size=14, bold=True, color=GRAPH, align=PP_ALIGN.CENTER)
    textbox(slide, 6.55, 1.68, 3.2, 0.55,
            "Multi-hop · Explainable\nAudit trail · Low hallucination",
            size=9, color=GRAPH, align=PP_ALIGN.CENTER)

    # Traditional RAG path (bottom)
    rrect(slide, 6.55, 2.88, 3.2, 1.15, fill=RGBColor(0x03, 0x0E, 0x28), line=TRAD, lw=1.5)
    textbox(slide, 6.55, 2.94, 3.2, 0.45,
            "Traditional RAG", size=14, bold=True, color=TRAD, align=PP_ALIGN.CENTER)
    textbox(slide, 6.55, 3.36, 3.2, 0.55,
            "Fast · Simple · Tier 1\nSub-2.3s latency",
            size=9, color=TRAD, align=PP_ALIGN.CENTER)

    # Answer output
    textbox(slide, 9.75, 2.42, 0.4, 0.4, "▶", size=14, color=DGREY, align=PP_ALIGN.CENTER)
    rrect(slide, 10.15, 2.2, 2.7, 0.85, fill=CARD, line=WHITE, lw=0.8)
    textbox(slide, 10.15, 2.2, 2.7, 0.85,
            "Answer + Audit Trail", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Benefit cards below
    benefits = [
        ("94%", "of Graph RAG's accuracy benefit captured", GRAPH),
        ("~2.5s", "average latency across the full query mix", TRAD),
        ("100%", "compliance queries routed to Graph RAG always", GOLD),
    ]
    for i, (num, label, color) in enumerate(benefits):
        lx = 0.45 + i * 4.3
        rrect(slide, lx, 4.55, 4.0, 2.2, fill=CARD, line=BORDER)
        rect(slide, lx, 4.55, 4.0, 0.065, fill=color)
        textbox(slide, lx, 4.75, 4.0, 0.88,
                num, size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(slide, lx + 0.15, 5.62, 3.7, 0.85,
                label, size=11, color=LGREY, align=PP_ALIGN.CENTER)

    add_fade_transition(slide)
    return slide


# ── SLIDE 12: GDPR & COMPLIANCE ───────────────────────────────────────────────
def s_gdpr(prs):
    slide = blank_slide(prs)
    slide_header(slide, "GDPR and Regulatory Compliance")

    rows_data = [
        ("Data Minimisation",    "Compliant — top-k chunks only",       "Requires scope-limited traversal",   GRAPH),
        ("MiFID II Explainability","Limited — no reasoning path",        "Full traversal path logged per query",GRAPH),
        ("Right to Erasure",     "Simple — delete chunk from vector store","Complex — cascading node deletion", GOLD),
        ("Audit Trail",          "Weak — no explicit reasoning log",     "Strong — node/edge provenance",       GRAPH),
        ("Access Control",       "Collection-level RBAC",               "Edge/path-level RBAC required",       GOLD),
        ("Hallucination Risk",   "Medium (21%)",                        "Low (9%) — 57% reduction",            GRAPH),
    ]

    headers = ["Compliance Dimension", "Traditional RAG", "Graph RAG"]
    col_widths = [3.3, 4.2, 4.2]
    col_starts = [0.45, 3.85, 8.15]
    row_h = 0.82

    for j, (hdr, lx, cw) in enumerate(zip(headers, col_starts, col_widths)):
        rrect(slide, lx, 0.9, cw, 0.45, fill=NAVY, line=BLUE, lw=0.8)
        textbox(slide, lx, 0.9, cw, 0.45,
                hdr, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (dim, trad_v, graph_v, winner_color) in enumerate(rows_data):
        ty = 1.42 + i * row_h
        bg = CARD if i % 2 == 0 else RGBColor(0x0A, 0x12, 0x22)
        for j, (val, lx, cw) in enumerate(zip([dim, trad_v, graph_v], col_starts, col_widths)):
            rrect(slide, lx, ty, cw, row_h - 0.06, fill=bg, line=BORDER, lw=0.3)
            if j == 0:
                textbox(slide, lx + 0.12, ty + 0.12, cw - 0.15, row_h - 0.1,
                        val, size=10.5, bold=True, color=LGREY)
            elif j == 1:
                textbox(slide, lx + 0.1, ty + 0.1, cw - 0.15, row_h - 0.1,
                        val, size=10, color=LGREY)
            else:
                textbox(slide, lx + 0.1, ty + 0.1, cw - 0.15, row_h - 0.1,
                        val, size=10, bold=True, color=winner_color)

    rrect(slide, 0.45, 6.48, 12.43, 0.72, fill=RGBColor(0x03, 0x10, 0x28), line=BLUE, lw=1.0)
    textbox(slide, 0.65, 6.56, 12.0, 0.58,
            "Before production deployment: implement cascading deletion for Graph RAG, "
            "add edge/path-level RBAC, and restrict graph inference on sensitive attributes.",
            size=10.5, color=BLUE_L)

    add_fade_transition(slide)
    return slide


# ── SLIDE 13: ROADMAP ─────────────────────────────────────────────────────────
def s_roadmap(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Deployment Roadmap")

    phases = [
        ("Q3 2026", "Phase 1\nTraditional RAG",
         ["Deploy vector RAG for internal FAQ", "Collect live query logs",
          "Train complexity classifier dataset"],
         TRAD),
        ("Q4 2026", "Phase 2\nGraph RAG (Compliance)",
         ["Deploy Graph RAG for compliance officers", "MiFID II suitability queries",
          "AML EDD decision support"],
         GRAPH),
        ("Q1 2027", "Phase 3\nHybrid Routing",
         ["Launch BERT query classifier", "Unified API layer for both pipelines",
          "Measure live latency vs accuracy"],
         BLUE),
        ("Q2 2027", "Phase 4\nExpand Corpus",
         ["Spanish regulatory documents", "Banco de Espana circulars, CNMV",
          "DORA and EU AI Act analysis"],
         GOLD),
    ]

    card_w = 2.9
    for i, (quarter, title, bullets, color) in enumerate(phases):
        lx = 0.45 + i * 3.22
        # Timeline connector
        if i < 3:
            rect(slide, lx + card_w, 2.08, 0.32, 0.07, fill=DGREY)
            textbox(slide, lx + card_w + 0.05, 1.85, 0.22, 0.3,
                    "▶", size=9, color=DGREY, align=PP_ALIGN.CENTER)

        # Phase badge
        rrect(slide, lx + 0.4, 0.92, card_w - 0.8, 0.42,
              fill=RGBColor(max(color[0]-80,0), max(color[1]-80,0), max(color[2]-80,0)),
              line=color, lw=1.0)
        textbox(slide, lx + 0.4, 0.92, card_w - 0.8, 0.42,
                quarter, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Phase card
        rrect(slide, lx, 1.42, card_w, 5.35, fill=CARD, line=BORDER)
        rect(slide, lx, 1.42, card_w, 0.07, fill=color)
        textbox(slide, lx + 0.15, 1.55, card_w - 0.2, 0.82,
                title, size=13, bold=True, color=color)

        for j, bullet_text in enumerate(bullets):
            ty = 2.52 + j * 1.28
            rrect(slide, lx + 0.15, ty, card_w - 0.3, 1.1,
                  fill=RGBColor(0x07, 0x0D, 0x1A), line=BORDER)
            rect(slide, lx + 0.15, ty, 0.045, 1.1, fill=color)
            textbox(slide, lx + 0.28, ty + 0.15, card_w - 0.48, 0.82,
                    bullet_text, size=10, color=LGREY)

    add_fade_transition(slide)
    return slide


# ── SLIDE 14: THE CORPUS ──────────────────────────────────────────────────────
def s_corpus(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Research Methodology: Corpus and Evaluation Design")

    # Left: corpus
    rrect(slide, 0.45, 0.92, 6.15, 6.28, fill=CARD, line=BORDER)
    rect(slide, 0.45, 0.92, 6.15, 0.065, fill=BLUE)
    textbox(slide, 0.65, 1.0, 5.75, 0.42,
            "Corpus: 9 Documents, ~750 Pages", size=12, bold=True, color=BLUE)

    docs = [
        ("EBA AML/CFT Guidelines 2021",       "87 pp",   GRAPH),
        ("BIS Basel III Framework",            "63 pp",   TRAD),
        ("MiFID II Directive 2014/65/EU",      "90 pp",   TRAD),
        ("GDPR Regulation 2016/679/EU",        "54 pp",   TRAD),
        ("Banc Sabadell Pillar III 2025",      "210 pp",  BLUE),
        ("Banc Sabadell Annual Report 2025",   "~350 pp", BLUE),
        ("EBA AML sample text",                "15 pp",   DGREY),
        ("Basel III sample text",              "12 pp",   DGREY),
        ("GDPR Banking sample text",           "14 pp",   DGREY),
    ]
    for i, (name, pages, color) in enumerate(docs):
        ty = 1.52 + i * 0.6
        rect(slide, 0.65, ty + 0.08, 0.045, 0.42, fill=color)
        textbox(slide, 0.75, ty + 0.04, 4.5, 0.42,
                name, size=10, color=LGREY)
        textbox(slide, 5.25, ty + 0.04, 1.15, 0.42,
                pages, size=10, bold=True, color=color, align=PP_ALIGN.RIGHT)

    # Right: evaluation design
    rrect(slide, 6.78, 0.92, 6.1, 6.28, fill=CARD, line=BORDER)
    rect(slide, 6.78, 0.92, 6.1, 0.065, fill=GOLD)
    textbox(slide, 6.98, 1.0, 5.7, 0.42,
            "Evaluation: 75 Annotated Pairs", size=12, bold=True, color=GOLD)

    tiers_eval = [
        ("Tier 1", "Simple Factual",      "n=20", "Single-document lookups",             TRAD),
        ("Tier 2", "Single-hop",          "n=20", "Two linked entities",                  GRAPH),
        ("Tier 3", "Multi-hop Relational","n=20", "Three or more entity hops",            GRAPH),
        ("Tier 4", "Adversarial",         "n=15", "Designed to elicit hallucination",    GOLD),
    ]
    for i, (tier, name, n, desc, color) in enumerate(tiers_eval):
        ty = 1.52 + i * 1.35
        rrect(slide, 6.98, ty, 5.7, 1.18, fill=RGBColor(0x07, 0x0D, 0x1A), line=BORDER)
        rect(slide, 6.98, ty, 0.055, 1.18, fill=color)
        textbox(slide, 7.1, ty + 0.06, 2.0, 0.42,
                f"{tier}: {name}", size=11, bold=True, color=color)
        textbox(slide, 7.1, ty + 0.52, 2.0, 0.38,
                n, size=20, bold=True, color=WHITE)
        textbox(slide, 9.0, ty + 0.3, 2.55, 0.6,
                desc, size=10, color=LGREY)

    textbox(slide, 6.98, 7.0, 5.7, 0.42,
            "20% of pairs cross-validated by GPT-4o as second annotator.",
            size=9.5, color=DGREY, italic=True)

    add_fade_transition(slide)
    return slide


# ── SLIDE 15: THANK YOU ───────────────────────────────────────────────────────
def s_thankyou(prs):
    slide = blank_slide(prs)

    # Glow orbs
    glow = slide.shapes.add_shape(9, Inches(3.5), Inches(1.0), Inches(6.5), Inches(6.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x00, 0x30, 0x80)
    glow.line.fill.background()
    sp_pr = glow._element.find(qn('p:spPr'))
    sf = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if sf is not None:
        srgb = sf.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            a = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            a.set('val', '10000')

    rect(slide, 0, 0, 13.33, 0.065, fill=BLUE)
    rect(slide, 0, 7.43, 13.33, 0.07, fill=BLUE)

    logo(slide, 5.3, 0.5, 2.72)
    divider(slide, 3.5, 1.22, 6.33, color=RGBColor(0x1E, 0x3A, 0x6E))

    multitext(slide, 1.0, 1.38, 11.33, 1.5, [
        ("Thank You", 60, True, WHITE, PP_ALIGN.CENTER, 0),
        ("Questions?", 28, False, BLUE_L, PP_ALIGN.CENTER, 8),
    ])

    divider(slide, 2.0, 3.2, 9.33, color=BORDER)

    # Team cards
    team = [
        ("MAE", "Mohamed Aymen\nElmezouari", "Team Lead / AI Architecture"),
        ("AM",  "Abhay\nMathahalli",         "Knowledge Graph / Graph RAG"),
        ("WL",  "Warren\nLiu",               "Evaluation / Reporting"),
    ]
    for i, (initials, name, role) in enumerate(team):
        lx = 1.5 + i * 3.55
        circle(slide, lx + 0.5, 3.5, 1.1, fill=CARD)
        textbox(slide, lx + 0.5, 3.5, 1.1, 1.1,
                initials, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        textbox(slide, lx, 4.75, 2.1, 0.6,
                name, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, lx, 5.38, 2.1, 0.45,
                role, size=9, color=LGREY, align=PP_ALIGN.CENTER)

    divider(slide, 2.0, 6.05, 9.33, color=BORDER)

    multitext(slide, 1.0, 6.2, 11.33, 0.9, [
        ("ESADE Business School  |  Banc Sabadell Capstone 2026  |  June 2026",
         10, False, DGREY, PP_ALIGN.CENTER, 0),
        ("github.com/mohamedaymenelmezouari/rag-comparison-sabadell  |  "
         "rag-comparison-banc-sabadell.streamlit.app",
         9.5, False, RGBColor(0x1E, 0x3A, 0x6E), PP_ALIGN.CENTER, 5),
    ])

    add_fade_transition(slide)
    return slide


# ── MAIN ──────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    s_title(prs)
    s_agenda(prs)
    s_challenge(prs)
    s_rag(prs)
    s_trad_arch(prs)
    s_graph_arch(prs)
    s_wow(prs)
    s_results(prs)
    s_by_tier(prs)
    s_head2head(prs)
    s_recommendation(prs)
    s_gdpr(prs)
    s_roadmap(prs)
    s_corpus(prs)
    s_thankyou(prs)

    prs.save(str(OUT_PATH))
    size_kb = OUT_PATH.stat().st_size // 1024
    print(f"Saved {len(prs.slides)} slides → {OUT_PATH}  ({size_kb} KB)")


if __name__ == "__main__":
    build()
